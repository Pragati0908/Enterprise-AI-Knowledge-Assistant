from pptx import Presentation


class PPTXParser:

    @staticmethod
    def extract_text(file_path: str):

        presentation = Presentation(file_path)

        text = ""

        for slide in presentation.slides:

            for shape in slide.shapes:

                if hasattr(shape, "text"):

                    text += shape.text + "\n"

        return text